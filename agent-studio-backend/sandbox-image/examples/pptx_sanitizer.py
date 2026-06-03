"""
Single-node PPTX Sanitizer
===========================

Paste this script into a Code Executor node.  It performs the full
sanitization workflow in one node using multi-pause replay:

Run 1  →  ask(file)  pauses for upload
Run 2  →  ask(file) replays, text extracted, llm.complete() runs,
            ask(selection) pauses for user pick
Run 3  →  both ask() replay, replacement runs, output.file() emits result
"""

from agent_studio import output, llm
from pptx import Presentation
import json

# ── Step 1: Upload PPTX ─────────────────────────────────────────────
filepath = output.ask(
    "Upload your PowerPoint file",
    type="file",
    accept=".pptx",
)

# ── Step 2: Extract all text from slides ─────────────────────────────
prs = Presentation(filepath)
all_text = "\n".join(
    shape.text
    for slide in prs.slides
    for shape in slide.shapes
    if shape.has_text_frame
)

# ── Step 3: LLM identifies client / company / person names ──────────
raw = llm.complete(
    prompt=(
        "Extract ALL client names, company names, and person names "
        "from the following text.  Return ONLY a JSON object with a "
        '"names" key containing an array of unique strings.\n\n'
        f"{all_text}"
    ),
    model="bedrock.anthropic.claude-haiku-4-5",
    system_prompt="You are a precise entity extractor. Return valid JSON only.",
    output_schema={
        "type": "object",
        "properties": {
            "names": {
                "type": "array",
                "items": {"type": "string"},
            }
        },
        "required": ["names"],
    },
)
names = json.loads(raw)["names"]

if not names:
    output.data({"message": "No client or person names found in the presentation."})
    raise SystemExit(0)

# ── Step 4: User picks which name to KEEP ────────────────────────────
names_list = "\n".join(f"  • {n}" for n in names)
keep = output.ask(
    f"Found {len(names)} name(s) in the presentation:\n\n"
    f"{names_list}\n\n"
    "Select the client name to KEEP — all others will be replaced with [REDACTED].",
    type="selection",
    options=names,
)

# ── Step 5: Replace all other names with [REDACTED] ─────────────────
remove = [n for n in names if n != keep]

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for name in remove:
                        run.text = run.text.replace(name, "[REDACTED]")

# ── Step 6: Save and output the sanitized file with summary ─────────
prs.save("/workspace/sanitized.pptx")
output.file("/workspace/sanitized.pptx", display_name="Sanitized Presentation.pptx")

print(f"\nKept: {keep}")
print(f"Redacted ({len(remove)}): {', '.join(remove)}")
