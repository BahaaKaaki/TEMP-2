"""
PowerPoint Builder - python-pptx Implementation

Converts the slide JSON schema into actual .pptx files using python-pptx.
Each layout type has a dedicated render function that creates shapes, text boxes,
and styled containers matching the Strategy& brand from the HTML assets.

Coordinate system: inches (python-pptx uses EMU internally, we use Inches/Pt helpers).
Slide dimensions: 13.333" x 7.5" (960x540pt equivalent).
"""
import io
import logging
from typing import List, Optional, Dict, Any

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from services.powerpoint_schema import (
    Presentation,
    Slide,
    SlideTheme,
    LayoutType,
)

logger = logging.getLogger(__name__)


# =============================================================================
# HELPERS
# =============================================================================

def _rgb(hex_str: str) -> RGBColor:
    """Convert hex string (with or without #) to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_text_box(slide, text, x, y, w, h, font_name="Arial", font_size=12,
                  color="2d2d2d", bold=False, italic=False, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP):
    """Add a simple text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = _rgb(color)
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    
    # Set vertical anchor
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    txBox.text_frame.paragraphs[0].space_after = Pt(0)
    
    return txBox


def _add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.5):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_color)
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    
    return shape


def _add_circle(slide, x, y, size, fill_color):
    """Add a circle (ellipse) shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()
    return shape


# =============================================================================
# TEXT-FLOW HELPERS
# =============================================================================

def _estimate_lines(text: str, width_inches: float, font_size_pt: int,
                    font_name: str = "Arial") -> int:
    """Rough estimate of how many lines *text* will occupy in a text box."""
    avg_char_w = font_size_pt * 0.55
    if font_name.lower() == "georgia":
        avg_char_w *= 1.05
    chars_per_line = max(1, int(width_inches * 72 / avg_char_w))
    words = text.split()
    lines, cur = 1, 0
    for w in words:
        need = len(w) + (1 if cur else 0)
        if cur + need > chars_per_line:
            lines += 1
            cur = len(w)
        else:
            cur += need
    return lines


def _add_item_list(slide, items, x, y, w, h, theme,
                   bullet_style="numbered", font_size=11,
                   text_color=None):
    """Render a list of items as a **single** multi-paragraph text box.

    PowerPoint handles word-wrap and paragraph spacing internally so items
    never overlap regardless of text length.

    bullet_style: "numbered" | "arrow" | "bullet" | "dash" | None
    """
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    color = text_color or theme.text_main

    for i, item in enumerate(items):
        text = item if isinstance(item, str) else str(item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(3)

        if bullet_style:
            if bullet_style == "numbered":
                bul = f"{i + 1}."
            elif bullet_style == "arrow":
                bul = "\u2192"
            elif bullet_style == "dash":
                bul = "\u2013"
            else:
                bul = "\u2022"
            run_b = p.add_run()
            run_b.text = f"{bul} "
            run_b.font.name = "Arial"
            run_b.font.size = Pt(font_size)
            run_b.font.color.rgb = _rgb(theme.brand_color)
            run_b.font.bold = True

        run_t = p.add_run()
        run_t.text = text
        run_t.font.name = "Arial"
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = _rgb(color)

    return txBox


# =============================================================================
# MASTER ELEMENTS (shared across all slides)
# =============================================================================

FRAME_LEFT = 0.49
FRAME_WIDTH = 12.36
HEADLINE_Y = 0.42
HEADLINE_FONT = 24
SUBTITLE_FONT = 15
FOOTER_Y = 7.08
MIN_SUBTITLE_Y = 1.18
MIN_CONTENT_Y = 1.64


def _add_slide_master(slide, headline: str, subtitle: str, theme: SlideTheme,
                      footer_mid: str, page_num: int, total: int) -> float:
    """Add the common Strategy& slide elements and return the dynamic content_y.

    The subtitle and content Y are pushed down when long headlines wrap to
    multiple lines, preventing overlap.
    """
    lines = _estimate_lines(headline, FRAME_WIDTH, HEADLINE_FONT, theme.title_font)
    line_h = HEADLINE_FONT * 1.15 / 72
    headline_h = max(0.65, lines * line_h + 0.06)
    subtitle_y = max(MIN_SUBTITLE_Y, HEADLINE_Y + headline_h + 0.02)
    content_y = max(MIN_CONTENT_Y, subtitle_y + 0.30)

    _add_text_box(
        slide, headline,
        x=FRAME_LEFT, y=HEADLINE_Y, w=FRAME_WIDTH, h=headline_h,
        font_name=theme.title_font, font_size=HEADLINE_FONT,
        color=theme.text_main
    )

    _add_text_box(
        slide, subtitle,
        x=FRAME_LEFT, y=subtitle_y, w=FRAME_WIDTH, h=0.3,
        font_name=theme.body_font, font_size=SUBTITLE_FONT,
        color=theme.brand_color, bold=True
    )

    _add_text_box(
        slide, "Strategy&",
        x=FRAME_LEFT, y=FOOTER_Y, w=2, h=0.25,
        font_size=10, color="888888"
    )
    _add_text_box(
        slide, footer_mid,
        x=5.0, y=FOOTER_Y, w=3.33, h=0.25,
        font_size=10, color="888888", align=PP_ALIGN.CENTER
    )
    _add_text_box(
        slide, f"{page_num} / {total}",
        x=11.35, y=FOOTER_Y, w=1.5, h=0.25,
        font_size=10, color="888888", align=PP_ALIGN.RIGHT
    )

    return content_y


# =============================================================================
# LAYOUT RENDERERS
# =============================================================================

def _render_title_content(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render title_content layout (bullet list) as a single text frame."""
    raw_items = content.get("items", [])
    if not raw_items:
        return

    available_h = FOOTER_Y - content_y - 0.3
    txBox = slide.shapes.add_textbox(
        Inches(FRAME_LEFT), Inches(content_y),
        Inches(FRAME_WIDTH), Inches(available_h),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(raw_items):
        text = item.get("text", "") if isinstance(item, dict) else str(item)
        bold_prefix = item.get("bold_prefix") if isinstance(item, dict) else None
        level = item.get("level", 0) if isinstance(item, dict) else 0
        bullet = "\u2022" if level == 0 else "\u2013"
        indent = "      " * level

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)

        run_b = p.add_run()
        run_b.text = f"{indent}{bullet} "
        run_b.font.name = "Arial"
        run_b.font.size = Pt(13)
        run_b.font.color.rgb = _rgb(theme.brand_color)
        run_b.font.bold = True

        if bold_prefix:
            run_bp = p.add_run()
            run_bp.text = f"{bold_prefix}: "
            run_bp.font.name = "Arial"
            run_bp.font.size = Pt(13)
            run_bp.font.color.rgb = _rgb(theme.text_main)
            run_bp.font.bold = True

        run_t = p.add_run()
        run_t.text = text
        run_t.font.name = "Arial"
        run_t.font.size = Pt(13)
        run_t.font.color.rgb = _rgb(theme.text_main)


def _render_card(slide, card: dict, x: float, y: float, w: float, h: float,
                 theme: SlideTheme, numbered: bool = True):
    """Render a single card container with red top bar.

    Items are placed in a single multi-paragraph text frame so PowerPoint
    handles word-wrap automatically and items never overlap.
    """
    _add_rect(slide, x, y, w, h, fill_color=theme.white, line_color=theme.grey_border)
    _add_rect(slide, x, y, w, 0.055, fill_color=theme.brand_color)

    cursor_y = y + 0.12
    text_w = w - 0.3
    tag = card.get("tag", "")
    if tag:
        _add_text_box(
            slide, tag.upper(),
            x=x + 0.15, y=cursor_y, w=text_w, h=0.2,
            font_size=10, color=theme.brand_color, bold=True
        )
        cursor_y += 0.22

    title = card.get("title", "")
    title_lines = _estimate_lines(title, text_w, 13)
    title_h = max(0.28, title_lines * 0.22)
    _add_text_box(
        slide, title,
        x=x + 0.15, y=cursor_y, w=text_w, h=title_h,
        font_size=13, color=theme.text_main, bold=True
    )
    cursor_y += title_h + 0.04

    desc = card.get("description", "")
    if desc:
        desc_lines = _estimate_lines(desc, text_w, 11)
        desc_h = max(0.25, desc_lines * 0.18)
        _add_text_box(
            slide, desc,
            x=x + 0.15, y=cursor_y, w=text_w, h=desc_h,
            font_size=11, color=theme.text_light, italic=True
        )
        cursor_y += desc_h + 0.04

    footer = card.get("footer", "")
    footer_text = footer.upper() if footer else ""
    footer_lines = _estimate_lines(footer_text, text_w, 9) if footer_text else 0
    footer_h = max(0.40, footer_lines * 0.16 + 0.12) if footer else 0.0
    items = card.get("items", [])
    if items:
        items_h = y + h - cursor_y - footer_h - 0.06
        _add_item_list(
            slide, items,
            x=x + 0.15, y=cursor_y, w=text_w, h=max(items_h, 0.5),
            theme=theme,
            bullet_style="numbered" if numbered else "bullet",
            font_size=11,
        )

    if footer:
        fY = y + h - footer_h
        _add_rect(slide, x, fY, w, footer_h, fill_color=theme.brand_color_light)
        _add_text_box(
            slide, footer_text,
            x=x + 0.15, y=fY + 0.04, w=text_w, h=footer_h - 0.08,
            font_size=9, color=theme.brand_color_dark, bold=True
        )


def _render_card_grid(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render two_column or three_column card layouts."""
    cards = content.get("cards", [])
    if not cards:
        return

    n = len(cards)
    gap = 0.20
    card_w = (FRAME_WIDTH - gap * (n - 1)) / n
    card_h = FOOTER_Y - content_y - 0.25
    start_x = FRAME_LEFT

    for i, card in enumerate(cards):
        cx = start_x + i * (card_w + gap)
        _render_card(slide, card, cx, content_y, card_w, card_h, theme)


def _render_spotlight(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render spotlight layout: 1 large + sidebar cards."""
    main = content.get("main_card", {})
    sidebar = content.get("sidebar_cards", [])
    stats = content.get("stats", [])

    sX = FRAME_LEFT
    sY = content_y
    main_w = 7.5
    side_w = 4.6
    gap = 0.26
    h = FOOTER_Y - content_y - 0.25
    
    # Main card
    _render_card(slide, main, sX, sY, main_w, h, theme)
    
    # Stats row inside main card
    if stats:
        stat_y = sY + h - 0.85
        # Divider line
        _add_rect(slide, sX + 0.15, stat_y, main_w - 0.3, 0.01, fill_color=theme.grey_light)
        for i, st in enumerate(stats):
            stX = sX + 0.3 + (i * 2.35)
            _add_text_box(
                slide, st.get("value", ""),
                x=stX, y=stat_y + 0.12, w=2, h=0.4,
                font_size=24, color=theme.brand_color, bold=True, align=PP_ALIGN.CENTER
            )
            _add_text_box(
                slide, st.get("label", ""),
                x=stX, y=stat_y + 0.52, w=2, h=0.2,
                font_size=12, color=theme.text_light, align=PP_ALIGN.CENTER
            )
    
    # Sidebar cards
    sbX = sX + main_w + gap
    sb_h = (h - gap) / max(len(sidebar), 1)
    for i, sb in enumerate(sidebar):
        sb_y = sY + i * (sb_h + gap)
        _render_card(slide, sb, sbX, sb_y, side_w, sb_h, theme)


def _render_approach(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render approach phases (3/4/5 step).

    Substeps are rendered as a single multi-paragraph text box per phase to
    avoid overlap when step descriptions wrap.
    """
    phases = content.get("phases", [])
    if not phases:
        return

    n = len(phases)
    has_crosscut = content.get("crosscut") is not None
    crosscut_reserve = 0.85 if has_crosscut else 0.0
    available_h = FOOTER_Y - content_y - 0.25 - crosscut_reserve

    gap = 0.25 if n <= 3 else (0.18 if n == 4 else 0.14)
    phase_w = (FRAME_WIDTH - gap * (n - 1)) / n
    phase_h = available_h
    phase_y = content_y

    num_size = 0.55 if n <= 3 else (0.48 if n == 4 else 0.40)
    num_font = 20 if n <= 3 else (18 if n == 4 else 16)
    title_font = 16 if n <= 3 else (15 if n == 4 else 14)

    for i, phase in enumerate(phases):
        pX = FRAME_LEFT + i * (phase_w + gap)

        _add_rect(slide, pX, phase_y, phase_w, phase_h,
                  fill_color=theme.white, line_color=theme.grey_border)
        _add_rect(slide, pX, phase_y, phase_w, 0.07, fill_color=theme.brand_color)

        _add_circle(slide, pX + 0.2, phase_y + 0.25, num_size, theme.brand_color)
        _add_text_box(
            slide, str(phase.get("number", i + 1)),
            x=pX + 0.2, y=phase_y + 0.32, w=num_size, h=num_size - 0.13,
            font_size=num_font, color=theme.white, bold=True, align=PP_ALIGN.CENTER
        )

        _add_text_box(
            slide, phase.get("title", ""),
            x=pX + 0.2 + num_size + 0.1, y=phase_y + 0.28,
            w=phase_w - num_size - 0.5, h=0.35,
            font_size=title_font, color=theme.text_main, bold=True
        )

        duration = phase.get("duration", "")
        if duration:
            _add_text_box(
                slide, duration,
                x=pX + 0.2 + num_size + 0.1, y=phase_y + 0.58,
                w=phase_w - num_size - 0.5, h=0.25,
                font_size=12, color=theme.text_light
            )

        substeps = phase.get("substeps", [])
        if substeps:
            ss_y = phase_y + 1.0
            ss_h = phase_h - 1.1
            _add_item_list(
                slide, substeps,
                x=pX + 0.15, y=ss_y, w=phase_w - 0.3, h=max(ss_h, 0.5),
                theme=theme,
                bullet_style="numbered",
                font_size=11,
            )

    crosscut = content.get("crosscut")
    if crosscut:
        cc_y = phase_y + phase_h + 0.2
        _add_rect(slide, FRAME_LEFT, cc_y, FRAME_WIDTH, 0.65, fill_color=theme.brand_color)
        _add_text_box(
            slide, crosscut.get("title", ""),
            x=0.95, y=cc_y + 0.18, w=3.2, h=0.3,
            font_size=14, color=theme.white, bold=True
        )
        items_text = "  \u2022  ".join(crosscut.get("items", []))
        _add_text_box(
            slide, items_text,
            x=4.3, y=cc_y + 0.2, w=8.3, h=0.28,
            font_size=12, color=theme.white
        )


def _render_metrics_dashboard(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render metrics dashboard: KPI boxes + insight panel."""
    metrics = content.get("metrics", [])
    insights = content.get("insights", [])

    n = len(metrics)
    if n == 0:
        return

    mW = (FRAME_WIDTH - 0.16 * (n - 1)) / n
    mH = 1.55
    sX = FRAME_LEFT
    sY = content_y
    
    for i, m in enumerate(metrics):
        mX = sX + i * (mW + 0.16)
        _add_rect(slide, mX, sY, mW, mH, fill_color=theme.white, line_color=theme.grey_border)
        _add_rect(slide, mX, sY, mW, 0.055, fill_color=theme.brand_color)
        
        _add_text_box(
            slide, m.get("value", ""),
            x=mX, y=sY + 0.2, w=mW, h=0.5,
            font_size=28, color=theme.brand_color, bold=True, align=PP_ALIGN.CENTER
        )
        _add_text_box(
            slide, m.get("label", ""),
            x=mX, y=sY + 0.7, w=mW, h=0.25,
            font_size=12, color=theme.text_light, align=PP_ALIGN.CENTER
        )
        
        change = m.get("change", "")
        if change:
            status = m.get("status", "neutral")
            change_color = theme.green if status == "up" else (theme.red_bright if status == "down" else theme.text_light)
            change_bg = theme.green_light if status == "up" else theme.grey_light
            _add_rect(slide, mX + 0.5, sY + 1.05, mW - 1.0, 0.3, fill_color=change_bg)
            _add_text_box(
                slide, change,
                x=mX + 0.5, y=sY + 1.07, w=mW - 1.0, h=0.26,
                font_size=12, color=change_color, align=PP_ALIGN.CENTER
            )
    
    if insights:
        pY = sY + mH + 0.2
        pH = 3.0
        _add_rect(slide, sX, pY, 12.36, pH, fill_color=theme.white, line_color=theme.grey_border)
        
        col_w = 12.36 / max(len(insights), 1)
        for i, ins in enumerate(insights):
            cX = sX + 0.2 + i * col_w
            
            # Divider line between columns
            if i > 0:
                _add_rect(slide, sX + i * col_w, pY + 0.15, 0.01, pH - 0.3,
                          fill_color=theme.grey_light)
            
            _add_text_box(
                slide, ins.get("title", ""),
                x=cX, y=pY + 0.2, w=col_w - 0.4, h=0.35,
                font_size=13, color=theme.text_main, bold=True
            )

            ins_items = ins.get("items", [])
            if ins_items:
                _add_item_list(
                    slide, ins_items,
                    x=cX, y=pY + 0.55, w=col_w - 0.4, h=pH - 0.7,
                    theme=theme,
                    bullet_style="bullet",
                    font_size=11,
                    text_color=theme.text_light,
                )


def _render_timeline(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render timeline with phase cards.

    Items inside each card use a single multi-paragraph text box so
    wrapped text never overlaps.
    """
    phases = content.get("phases", [])
    if not phases:
        return

    n = len(phases)
    sX = FRAME_LEFT
    sY = content_y
    track_w = FRAME_WIDTH - 1.0

    _add_rect(slide, sX + 0.5, sY + 0.35, track_w, 0.04, fill_color=theme.grey_border)

    for i, p in enumerate(phases):
        dot_x = sX + 0.6 + i * (track_w / max(n - 1, 1)) - 0.1
        dot_color = theme.brand_color if p.get("status") != "future" else theme.white
        shape = _add_circle(slide, dot_x, sY + 0.27, 0.2, dot_color)
        if p.get("status") == "future":
            shape.line.color.rgb = _rgb(theme.brand_color)
            shape.line.width = Pt(2)

    cY = sY + 0.8
    cW = (FRAME_WIDTH - 0.22 * (n - 1)) / n
    cH = FOOTER_Y - cY - 0.25

    for i, p in enumerate(phases):
        cX = sX + i * (cW + 0.22)
        _add_rect(slide, cX, cY, cW, cH, fill_color=theme.white, line_color=theme.grey_border)

        status = p.get("status", "current")
        header_color = theme.text_light if status == "past" else (
            "b34747" if status == "future" else theme.brand_color
        )
        _add_rect(slide, cX, cY, cW, 0.5, fill_color=header_color)
        _add_text_box(
            slide, p.get("phase_label", ""),
            x=cX + 0.12, y=cY + 0.1, w=1.5, h=0.3,
            font_size=12, color=theme.white, bold=True
        )
        date_range = p.get("date_range", "")
        if date_range:
            _add_text_box(
                slide, date_range,
                x=cX + cW - 1.5, y=cY + 0.12, w=1.3, h=0.25,
                font_size=12, color=theme.white, align=PP_ALIGN.RIGHT
            )

        _add_text_box(
            slide, p.get("title", ""),
            x=cX + 0.15, y=cY + 0.6, w=cW - 0.3, h=0.35,
            font_size=13, color=theme.text_main, bold=True
        )

        items = p.get("items", [])
        if items:
            items_y = cY + 1.0
            items_h = cH - 1.1
            _add_item_list(
                slide, items,
                x=cX + 0.15, y=items_y, w=cW - 0.3, h=max(items_h, 0.5),
                theme=theme,
                bullet_style="arrow",
                font_size=10,
                text_color=theme.text_light,
            )


def _render_comparison_table(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render comparison table."""
    headers = content.get("headers", [])
    rows = content.get("rows", [])
    recommended = content.get("recommended_column")
    recommendation = content.get("recommendation", "")

    if not headers:
        return

    sX = FRAME_LEFT
    sY = content_y
    row_h = 0.6
    n_cols = len(headers)
    col_w = 12.36 / n_cols
    
    # Header row
    for i, h in enumerate(headers):
        hX = sX + i * col_w
        bg = theme.grey_light if i == 0 else (theme.green if i == recommended else theme.brand_color)
        text_color = theme.text_main if i == 0 else theme.white
        _add_rect(slide, hX, sY, col_w, row_h, fill_color=bg, line_color=bg)
        _add_text_box(
            slide, h,
            x=hX + 0.15, y=sY + 0.15, w=col_w - 0.3, h=0.32,
            font_size=12, color=text_color, bold=True
        )
    
    # Data rows
    for ri, row in enumerate(rows):
        rY = sY + row_h + ri * row_h
        for ci, cell in enumerate(row):
            cX = sX + ci * col_w
            bg = theme.grey_light if ci == 0 else (
                "fafafa" if ri % 2 == 1 else theme.white
            )
            _add_rect(slide, cX, rY, col_w, row_h, fill_color=bg, line_color=theme.grey_border)
            is_bold = ci == 0 or (ci == recommended and ri == 3)
            _add_text_box(
                slide, cell,
                x=cX + 0.15, y=rY + 0.15, w=col_w - 0.3, h=0.32,
                font_size=12, color=theme.text_main, bold=is_bold
            )
    
    if recommendation:
        fY = sY + row_h + len(rows) * row_h + 0.2
        _add_rect(slide, sX, fY, 12.36, 0.55, fill_color=theme.brand_color_light)
        _add_text_box(
            slide, f"Recommendation: {recommendation}",
            x=sX + 0.3, y=fY + 0.12, w=11.77, h=0.32,
            font_size=12, color=theme.brand_color_dark
        )


def _render_call_to_action(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render CTA: banner + action cards."""
    sX = FRAME_LEFT
    sY = content_y

    _add_rect(slide, sX, sY, FRAME_WIDTH, 1.15, fill_color=theme.grey_light, line_color=theme.grey_border)
    _add_rect(slide, sX + 0.25, sY + 0.2, 0.85, 0.75, fill_color=theme.brand_color)
    _add_text_box(
        slide, "!",
        x=sX + 0.25, y=sY + 0.32, w=0.85, h=0.5,
        font_size=28, color=theme.white, align=PP_ALIGN.CENTER, bold=True
    )
    _add_text_box(
        slide, content.get("banner_headline", ""),
        x=sX + 1.3, y=sY + 0.22, w=10.5, h=0.4,
        font_size=15, color=theme.text_main, bold=True
    )
    subtext = content.get("banner_subtext", "")
    if subtext:
        _add_text_box(
            slide, subtext,
            x=sX + 1.3, y=sY + 0.65, w=10.5, h=0.35,
            font_size=12, color=theme.text_light
        )
    
    actions = content.get("actions", [])
    if not actions:
        return

    aY = sY + 1.4
    n = len(actions)
    gap = 0.18
    aW = (FRAME_WIDTH - gap * (n - 1)) / n
    aH = FOOTER_Y - aY - 0.25
    
    footer_zone = 0.55
    for i, a in enumerate(actions):
        aX = sX + i * (aW + gap)
        _add_rect(slide, aX, aY, aW, aH, fill_color=theme.white, line_color=theme.grey_border)
        _add_rect(slide, aX, aY, 0.055, aH, fill_color=theme.brand_color)

        _add_circle(slide, aX + 0.2, aY + 0.2, 0.42, theme.brand_color)
        _add_text_box(
            slide, str(a.get("number", i + 1)),
            x=aX + 0.2, y=aY + 0.27, w=0.42, h=0.3,
            font_size=14, color=theme.white, bold=True, align=PP_ALIGN.CENTER
        )

        title_text = a.get("title", "")
        title_lines = _estimate_lines(title_text, aW - 0.92, 13)
        title_h = max(0.30, title_lines * 0.22)
        _add_text_box(
            slide, title_text,
            x=aX + 0.72, y=aY + 0.25, w=aW - 0.92, h=title_h,
            font_size=13, color=theme.text_main, bold=True
        )

        desc_y = aY + 0.25 + title_h + 0.10
        desc_h = max(0.5, aH - (desc_y - aY) - footer_zone - 0.06)
        _add_text_box(
            slide, a.get("description", ""),
            x=aX + 0.2, y=desc_y, w=aW - 0.4, h=desc_h,
            font_size=10, color=theme.text_light
        )

        _add_rect(slide, aX + 0.15, aY + aH - footer_zone, aW - 0.3, 0.01,
                  fill_color=theme.grey_light)

        owner = a.get("owner", "")
        if owner:
            _add_text_box(
                slide, f"Owner: {owner}",
                x=aX + 0.2, y=aY + aH - footer_zone + 0.08, w=2, h=0.25,
                font_size=10, color=theme.text_light
            )

        due = a.get("due_date", "")
        if due:
            _add_rect(slide, aX + aW - 1.2, aY + aH - footer_zone + 0.05, 1.0, 0.30,
                      fill_color=theme.brand_color_light)
            _add_text_box(
                slide, due,
                x=aX + aW - 1.2, y=aY + aH - footer_zone + 0.08, w=1.0, h=0.25,
                font_size=10, color=theme.brand_color, bold=True, align=PP_ALIGN.CENTER
            )


def _render_kpi_hero(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render KPI hero: large KPI left + detail cards right.

    Matches HTML layout: left panel flex 0 0 280px, right cards flex 1,
    gap 16px.  At 960px canvas -> 13.333in slide the mapping is:
      280px -> 3.89in, 16px gap -> 0.22in, remaining -> 8.25in
    """
    left_x = FRAME_LEFT
    left_y = content_y
    left_w = 3.89
    left_h = FOOTER_Y - content_y - 0.25
    gap = 0.22
    
    # Left panel
    _add_rect(slide, left_x, left_y, left_w, left_h,
              fill_color=theme.white, line_color=theme.grey_border)
    _add_rect(slide, left_x, left_y, left_w, 0.06, fill_color=theme.brand_color)
    
    pill = content.get("pill_label", "")
    pill_offset = 0.0
    if pill:
        _add_text_box(
            slide, pill.upper(),
            x=left_x + 0.25, y=left_y + 0.2, w=left_w - 0.5, h=0.25,
            font_size=9, color=theme.brand_color_dark, bold=True
        )
        pill_offset = 0.35
    
    _add_text_box(
        slide, content.get("kpi_value", ""),
        x=left_x + 0.25, y=left_y + 0.2 + pill_offset, w=left_w - 0.5, h=0.7,
        font_size=36, color=theme.brand_color, bold=True
    )
    _add_text_box(
        slide, content.get("kpi_label", ""),
        x=left_x + 0.25, y=left_y + 0.95 + pill_offset, w=left_w - 0.5, h=0.4,
        font_size=12, color=theme.text_main, bold=True
    )
    
    subnote = content.get("kpi_subnote", "")
    if subnote:
        _add_text_box(
            slide, subnote,
            x=left_x + 0.25, y=left_y + 1.4 + pill_offset, w=left_w - 0.5, h=0.8,
            font_size=10, color=theme.text_light
        )
    
    footnote = content.get("kpi_footnote", "")
    if footnote:
        _add_text_box(
            slide, footnote,
            x=left_x + 0.25, y=left_y + left_h - 0.6, w=left_w - 0.5, h=0.5,
            font_size=9, color=theme.text_light
        )
    
    cards = content.get("cards", [])
    card_x = left_x + left_w + gap
    card_w = left_x + FRAME_WIDTH - card_x
    n_cards = max(len(cards), 1)
    card_gap = 0.11
    card_h = (left_h - card_gap * (n_cards - 1)) / n_cards

    for i, card in enumerate(cards):
        cy = left_y + i * (card_h + card_gap)

        _add_rect(slide, card_x, cy, card_w, card_h,
                  fill_color=theme.white, line_color=theme.grey_border)
        _add_rect(slide, card_x, cy, 0.06, card_h, fill_color=theme.brand_color)

        icon = card.get("icon", str(i + 1))
        _add_circle(slide, card_x + 0.18, cy + 0.14, 0.28, theme.brand_color_light)
        _add_text_box(
            slide, icon,
            x=card_x + 0.18, y=cy + 0.16, w=0.28, h=0.24,
            font_size=9, color=theme.brand_color_dark, bold=True, align=PP_ALIGN.CENTER
        )

        _add_text_box(
            slide, card.get("title", ""),
            x=card_x + 0.55, y=cy + 0.14, w=card_w - 0.75, h=0.28,
            font_size=12, color=theme.text_main, bold=True
        )

        items = card.get("items", [])
        if items:
            items_y = cy + 0.45
            items_h = card_h - 0.55
            _add_item_list(
                slide, items,
                x=card_x + 0.18, y=items_y, w=card_w - 0.36, h=max(items_h, 0.3),
                theme=theme,
                bullet_style="arrow",
                font_size=10,
                text_color=theme.text_light,
            )


# =============================================================================
# CV LAYOUT RENDERERS
# =============================================================================

def _render_cv_team_summary(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render a team overview slide with profile cards in a grid (up to 6)."""
    profiles = content.get("profiles", [])
    if not profiles:
        return

    n = len(profiles)
    cols = min(n, 3)
    rows = 1 if n <= 3 else 2
    gap = 0.18
    card_w = (FRAME_WIDTH - gap * (cols - 1)) / cols
    total_h = FOOTER_Y - content_y - 0.25
    card_h = (total_h - gap * (rows - 1)) / rows

    for idx, profile in enumerate(profiles):
        row = idx // cols
        col = idx % cols
        cx = FRAME_LEFT + col * (card_w + gap)
        cy = content_y + row * (card_h + gap)

        _add_rect(slide, cx, cy, card_w, card_h,
                  fill_color=theme.white, line_color=theme.grey_border)
        _add_rect(slide, cx, cy, card_w, 0.055, fill_color=theme.brand_color)

        name = profile.get("name", "")
        level = profile.get("level", "")
        city = profile.get("city", "")
        meta = " | ".join(filter(None, [level, city]))

        cursor = cy + 0.14
        _add_text_box(slide, name,
                      x=cx + 0.15, y=cursor, w=card_w - 0.3, h=0.28,
                      font_size=13, color=theme.text_main, bold=True)
        cursor += 0.28
        if meta:
            _add_text_box(slide, meta,
                          x=cx + 0.15, y=cursor, w=card_w - 0.3, h=0.22,
                          font_size=10, color=theme.text_light, italic=True)
            cursor += 0.24

        _add_rect(slide, cx + 0.15, cursor, card_w - 0.3, 0.01,
                  fill_color=theme.grey_border)
        cursor += 0.08

        summary_lines = profile.get("summary_lines", [])
        if summary_lines:
            _add_text_box(slide, "Executive Summary",
                          x=cx + 0.15, y=cursor, w=card_w - 0.3, h=0.2,
                          font_size=9, color=theme.brand_color, bold=True)
            cursor += 0.2
            section_h = (cy + card_h - cursor) * 0.45
            _add_item_list(slide, summary_lines,
                           x=cx + 0.15, y=cursor, w=card_w - 0.3,
                           h=max(section_h, 0.3), theme=theme,
                           bullet_style="dash", font_size=9,
                           text_color=theme.text_light)
            cursor += section_h + 0.06

        rel_lines = profile.get("relevant_experience_lines", [])
        if rel_lines:
            _add_text_box(slide, "Relevant Experience",
                          x=cx + 0.15, y=cursor, w=card_w - 0.3, h=0.2,
                          font_size=9, color=theme.brand_color, bold=True)
            cursor += 0.2
            remaining = max(cy + card_h - cursor - 0.1, 0.3)
            _add_item_list(slide, rel_lines,
                           x=cx + 0.15, y=cursor, w=card_w - 0.3,
                           h=remaining, theme=theme,
                           bullet_style="dash", font_size=9,
                           text_color=theme.text_light)


def _render_cv_individual(slide, content: dict, theme: SlideTheme, content_y: float):
    """Render a two-column individual CV slide (profile left, projects right)."""
    left_w = 5.2
    right_w = FRAME_WIDTH - left_w - 0.25
    right_x = FRAME_LEFT + left_w + 0.25
    avail_h = FOOTER_Y - content_y - 0.25

    name = content.get("name", "")
    level = content.get("level", "")
    city = content.get("city", "")
    left_col = content.get("left_column", {})
    projects = content.get("projects", [])

    # -- Left column background --
    _add_rect(slide, FRAME_LEFT, content_y, left_w, avail_h,
              fill_color=theme.white, line_color=theme.grey_border)
    _add_rect(slide, FRAME_LEFT, content_y, left_w, 0.055,
              fill_color=theme.brand_color)

    cursor = content_y + 0.14
    _add_text_box(slide, name,
                  x=FRAME_LEFT + 0.2, y=cursor, w=left_w - 0.4, h=0.3,
                  font_size=14, color=theme.text_main, bold=True)
    cursor += 0.32
    meta = " | ".join(filter(None, [level, city]))
    if meta:
        _add_text_box(slide, meta,
                      x=FRAME_LEFT + 0.2, y=cursor, w=left_w - 0.4, h=0.22,
                      font_size=10, color=theme.text_light, italic=True)
        cursor += 0.26
    _add_rect(slide, FRAME_LEFT + 0.2, cursor, left_w - 0.4, 0.01,
              fill_color=theme.grey_border)
    cursor += 0.12

    sections = [
        ("Executive Summary", left_col.get("executive_summary", [])),
        ("Relevant Experience", left_col.get("relevant_experience", [])),
        ("Prior Experience", left_col.get("prior_experience", [])),
        ("Education & Languages", left_col.get("education_and_languages", [])),
    ]
    non_empty = [s for s in sections if s[1]]
    if non_empty:
        section_h = (content_y + avail_h - cursor - 0.1) / len(non_empty)
        for title, items in non_empty:
            _add_text_box(slide, title,
                          x=FRAME_LEFT + 0.2, y=cursor, w=left_w - 0.4, h=0.2,
                          font_size=9, color=theme.brand_color, bold=True)
            cursor += 0.2
            _add_item_list(slide, items,
                           x=FRAME_LEFT + 0.2, y=cursor, w=left_w - 0.4,
                           h=max(section_h - 0.28, 0.25), theme=theme,
                           bullet_style="dash", font_size=9,
                           text_color=theme.text_main)
            cursor += section_h - 0.2

    # -- Right column (projects) --
    _add_rect(slide, right_x, content_y, right_w, avail_h,
              fill_color=theme.white, line_color=theme.grey_border)
    _add_rect(slide, right_x, content_y, right_w, 0.055,
              fill_color=theme.brand_color)

    _add_text_box(slide, "Case Studies / Projects",
                  x=right_x + 0.2, y=content_y + 0.14, w=right_w - 0.4, h=0.25,
                  font_size=11, color=theme.brand_color, bold=True)

    proj_cursor = content_y + 0.48
    if projects:
        proj_h = (content_y + avail_h - proj_cursor - 0.1) / len(projects)
        for proj in projects:
            _add_text_box(slide, proj.get("title", ""),
                          x=right_x + 0.2, y=proj_cursor, w=right_w - 0.4, h=0.22,
                          font_size=10, color=theme.text_main, bold=True)
            proj_cursor += 0.22
            bullets = proj.get("bullets", [])
            if bullets:
                _add_item_list(slide, bullets,
                               x=right_x + 0.2, y=proj_cursor,
                               w=right_w - 0.4,
                               h=max(proj_h - 0.32, 0.25), theme=theme,
                               bullet_style="dash", font_size=9,
                               text_color=theme.text_light)
                proj_cursor += proj_h - 0.22


# =============================================================================
# LAYOUT DISPATCHER
# =============================================================================

LAYOUT_RENDERERS = {
    LayoutType.TITLE_CONTENT: _render_title_content,
    LayoutType.TWO_COLUMN: _render_card_grid,
    LayoutType.THREE_COLUMN: _render_card_grid,
    LayoutType.SPOTLIGHT: _render_spotlight,
    LayoutType.APPROACH_3STEP: _render_approach,
    LayoutType.APPROACH_4STEP: _render_approach,
    LayoutType.APPROACH_5STEP: _render_approach,
    LayoutType.HORIZONTAL_APPROACH: _render_approach,
    LayoutType.METRICS_DASHBOARD: _render_metrics_dashboard,
    LayoutType.TIMELINE: _render_timeline,
    LayoutType.COMPARISON_TABLE: _render_comparison_table,
    LayoutType.CALL_TO_ACTION: _render_call_to_action,
    LayoutType.KPI_HERO: _render_kpi_hero,
    LayoutType.CV_TEAM_SUMMARY: _render_cv_team_summary,
    LayoutType.CV_INDIVIDUAL: _render_cv_individual,
}


# =============================================================================
# PUBLIC API
# =============================================================================

def build_pptx(presentation: Presentation) -> bytes:
    """
    Build a .pptx file from a Presentation object.
    
    Returns the file as bytes (ready for HTTP response / download).
    """
    theme = presentation.theme
    slides = presentation.slides
    total = len(slides)
    
    logger.debug("Building PPTX: '%s' (%d slides)", presentation.title, total)
    
    pptx = PptxPresentation()
    
    # Set slide dimensions
    pptx.slide_width = Inches(theme.slide_width)
    pptx.slide_height = Inches(theme.slide_height)
    
    # Use blank layout
    blank_layout = pptx.slide_layouts[6]  # Blank layout
    
    for i, slide_data in enumerate(slides):
        pptx_slide = pptx.slides.add_slide(blank_layout)
        
        # Get content as dict
        content = slide_data.content
        if hasattr(content, "model_dump"):
            content = content.model_dump()
        elif not isinstance(content, dict):
            content = dict(content) if content else {}
        
        content_y = _add_slide_master(
            pptx_slide,
            headline=slide_data.headline,
            subtitle=slide_data.subtitle,
            theme=theme,
            footer_mid=slide_data.subtitle,
            page_num=i + 1,
            total=total,
        )

        renderer = LAYOUT_RENDERERS.get(slide_data.layout)
        if renderer:
            try:
                renderer(pptx_slide, content, theme, content_y)
            except Exception as e:
                logger.error(
                    "Error rendering slide %d ('%s', layout=%s): %s",
                    i + 1, slide_data.headline[:40], slide_data.layout, e,
                    exc_info=True,
                )
                # Add error text box as fallback
                _add_text_box(
                    pptx_slide,
                    f"Error rendering this slide: {str(e)}",
                    x=0.48, y=3.5, w=12.37, h=0.5,
                    font_size=14, color="c62828"
                )
        else:
            logger.warning("No renderer for layout '%s', skipping content", slide_data.layout)
        
        # Add speaker notes
        if slide_data.speaker_notes:
            pptx_slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes
    
    # Save to bytes
    buffer = io.BytesIO()
    pptx.save(buffer)
    buffer.seek(0)
    
    logger.debug("PPTX built successfully: %d bytes", buffer.getbuffer().nbytes)
    return buffer.getvalue()
