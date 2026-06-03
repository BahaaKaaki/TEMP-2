"""Extract placeholders from a PPTX template.

Scans every text frame (shapes, tables, grouped shapes) for the supported
placeholder kinds:

- ``{{ field }}``  -- simple text substitution
- ``{{* field }}`` -- bullet-array (one paragraph per list item)
- ``{{+ field }}`` -- repeat-field (paragraph group repeats per array item)
- ``{{# name }}`` / ``{{/ name }}`` -- slide-loop markers
- ``{{@ name | count }}`` -- variant marker (select slide by array length)

Returns a structured list that the schema generator can consume.
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from enum import Enum
from typing import Dict, List, Optional, Set, Tuple

from pptx import Presentation

from .run_merger import merge_multiline_placeholders, paragraph_text


class PlaceholderKind(str, Enum):
    TEXT = "text"
    BULLET_ARRAY = "bullet_array"
    REPEAT_FIELD = "repeat_field"
    LOOP_START = "loop_start"
    LOOP_END = "loop_end"
    VARIANT = "variant"


@dataclass
class Placeholder:
    raw: str
    kind: PlaceholderKind
    path: str
    slide_index: int
    loop_context: Optional[str] = None
    description: Optional[str] = None


_RE_PLACEHOLDER = re.compile(
    r"\{\{\s*"
    r"(?P<prefix>[#*/+@]?)\s*"
    r"(?P<path>[a-zA-Z_][a-zA-Z0-9_.]*)"
    r"\s*(?:\|\s*(?P<desc>[^}]+?))?"
    r"\s*\}\}"
)


def _iter_text_frames(slide):
    """Yield every text_frame from a slide, including tables and nested groups."""

    def _walk(shapes):
        for shape in shapes:
            if shape.has_text_frame:
                yield shape.text_frame
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
            if hasattr(shape, "shapes"):
                yield from _walk(shape.shapes)

    yield from _walk(slide.shapes)


def extract_placeholders(
    pptx_path: Optional[str] = None,
    prs: Optional[Presentation] = None,
) -> List[Placeholder]:
    """Parse a PPTX and return all detected placeholders.

    Accepts either a file *pptx_path* or an already-loaded *prs*
    :class:`Presentation` object (useful when the presentation has been
    modified in memory, e.g. after variant selection).
    """
    if prs is None:
        if pptx_path is None:
            raise ValueError("Either pptx_path or prs must be provided")
        prs = Presentation(pptx_path)
    results: List[Placeholder] = []
    seen: Set[Tuple[str, int]] = set()

    for slide_idx, slide in enumerate(prs.slides):
        for tf in _iter_text_frames(slide):
            merge_multiline_placeholders(tf)
            for para in tf.paragraphs:
                full_text = paragraph_text(para)
                for m in _RE_PLACEHOLDER.finditer(full_text):
                    prefix = m.group("prefix")
                    path = m.group("path")
                    raw = m.group(0)

                    if prefix == "#":
                        kind = PlaceholderKind.LOOP_START
                    elif prefix == "/":
                        kind = PlaceholderKind.LOOP_END
                    elif prefix == "*":
                        kind = PlaceholderKind.BULLET_ARRAY
                    elif prefix == "+":
                        kind = PlaceholderKind.REPEAT_FIELD
                    elif prefix == "@":
                        kind = PlaceholderKind.VARIANT
                    else:
                        kind = PlaceholderKind.TEXT

                    key = (raw, slide_idx)
                    if key in seen:
                        continue
                    seen.add(key)

                    desc = (m.group("desc") or "").strip() or None

                    results.append(Placeholder(
                        raw=raw,
                        kind=kind,
                        path=path,
                        slide_index=slide_idx,
                        description=desc,
                    ))

    _assign_loop_contexts(results)
    return results


def _assign_loop_contexts(placeholders: List[Placeholder]) -> None:
    """Tag each placeholder inside a loop with its loop variable name.

    Uses slide-index ranges rather than shape iteration order so that a
    loop-end marker placed early in a slide (e.g. in a footer placeholder)
    does not prematurely close the context for content shapes on the same
    slide.
    """
    ranges: List[Tuple[str, int, int]] = []
    starts: Dict[str, int] = {}
    for ph in placeholders:
        if ph.kind == PlaceholderKind.LOOP_START:
            starts[ph.path] = ph.slide_index
        elif ph.kind == PlaceholderKind.LOOP_END:
            if ph.path in starts:
                ranges.append((ph.path, starts.pop(ph.path), ph.slide_index))

    for ph in placeholders:
        if ph.kind in (PlaceholderKind.LOOP_START, PlaceholderKind.LOOP_END):
            continue
        for name, s_start, s_end in ranges:
            if s_start <= ph.slide_index <= s_end:
                ph.loop_context = name
                break


def detect_loops(placeholders: List[Placeholder]) -> List[Dict]:
    """Return ``{name, start_slide, end_slide, fields}`` for each loop block.

    *fields* contains every ``item.*`` placeholder path found between
    the ``{{# }}`` and ``{{/ }}`` markers, with the ``item.`` prefix stripped.
    """
    starts: Dict[str, int] = {}
    ends: Dict[str, int] = {}

    for ph in placeholders:
        if ph.kind == PlaceholderKind.LOOP_START:
            starts[ph.path] = ph.slide_index
        elif ph.kind == PlaceholderKind.LOOP_END:
            ends[ph.path] = ph.slide_index

    loops = []
    for name in starts:
        s_start = starts[name]
        s_end = ends.get(name, s_start)
        loop_fields: List[str] = []
        for ph in placeholders:
            if ph.kind in (PlaceholderKind.LOOP_START, PlaceholderKind.LOOP_END):
                continue
            if not (s_start <= ph.slide_index <= s_end):
                continue
            item_prefix = "item."
            if ph.path.startswith(item_prefix):
                field_name = ph.path[len(item_prefix):]
            else:
                field_name = ph.path
            if field_name not in loop_fields:
                loop_fields.append(field_name)

        loops.append({
            "name": name,
            "start_slide": s_start,
            "end_slide": s_end,
            "fields": loop_fields,
        })
    return loops


def detect_repeat_groups(
    placeholders: List[Placeholder],
) -> List[Dict]:
    """Identify paragraph-level repeat groups from ``{{+ }}`` placeholders.

    A repeat group is a set of ``{{+ }}`` and ``{{* }}`` placeholders that
    share the same array root (path minus the leaf segment).  Returns a list
    of dicts::

        {
            "array_root": "person1.projects_left",
            "fields": [
                {"leaf": "title",   "kind": "repeat_field", "description": ...},
                {"leaf": "bullets", "kind": "bullet_array", "description": ...},
            ],
            "slide_index": 2,
        }
    """
    roots: Dict[str, Dict] = {}

    for ph in placeholders:
        if ph.kind == PlaceholderKind.REPEAT_FIELD:
            parts = ph.path.rsplit(".", 1)
            if len(parts) != 2:
                continue
            array_root, leaf = parts
            key = (array_root, ph.slide_index)
            if key not in roots:
                roots[key] = {
                    "array_root": array_root,
                    "fields": [],
                    "slide_index": ph.slide_index,
                }
            roots[key]["fields"].append({
                "leaf": leaf,
                "kind": ph.kind.value,
                "description": ph.description,
            })

    for ph in placeholders:
        if ph.kind != PlaceholderKind.BULLET_ARRAY:
            continue
        parts = ph.path.rsplit(".", 1)
        if len(parts) != 2:
            continue
        array_root, leaf = parts
        key = (array_root, ph.slide_index)
        if key in roots:
            already = {f["leaf"] for f in roots[key]["fields"]}
            if leaf not in already:
                roots[key]["fields"].append({
                    "leaf": leaf,
                    "kind": ph.kind.value,
                    "description": ph.description,
                })

    return list(roots.values())


def detect_variants(placeholders: List[Placeholder]) -> List[Dict]:
    """Identify variant groups from ``{{@ array_path | count }}`` markers.

    Returns a list of dicts, one per variant group::

        {
            "name": "people",
            "variants": {6: [0], 5: [1], 4: [2], 3: [3], 2: [4]},
            "min_count": 2,
            "max_count": 6,
        }

    *variants* maps each count to the list of slide indices tagged with
    that count.  Multi-slide variants are supported (multiple slides can
    share the same ``{{@ name | count }}`` marker).
    """
    groups: Dict[str, Dict[int, List[int]]] = {}

    for ph in placeholders:
        if ph.kind != PlaceholderKind.VARIANT:
            continue
        desc = (ph.description or "").strip()
        if not desc.isdigit():
            continue
        count = int(desc)
        if ph.path not in groups:
            groups[ph.path] = {}
        groups[ph.path].setdefault(count, []).append(ph.slide_index)

    results: List[Dict] = []
    for name, variants in groups.items():
        counts = sorted(variants.keys())
        results.append({
            "name": name,
            "variants": variants,
            "min_count": counts[0] if counts else 0,
            "max_count": counts[-1] if counts else 0,
        })
    return results


def summarise(placeholders: List[Placeholder]) -> str:
    """Human-readable summary of detected placeholders."""
    lines = ["Detected placeholders:", ""]
    for ph in placeholders:
        ctx = f"  (inside loop: {ph.loop_context})" if ph.loop_context else ""
        lines.append(
            f"  slide {ph.slide_index + 1} | {ph.kind.value:14s} | "
            f"{ph.raw}{ctx}"
        )
    return "\n".join(lines)
