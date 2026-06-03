"""Sanitize uploaded PPTX templates before storage.

Strips add-in metadata shapes (e.g. think-cell) and prunes orphaned
relationships and embedded parts that would cause "needs repair" errors
when slides are later cloned during template filling.

Called once at upload time so the fill path stays simple.
"""
from __future__ import annotations

import logging
import posixpath
import zipfile
from io import BytesIO
from typing import Set
from xml.etree import ElementTree as ET

from lxml import etree
from pptx import Presentation

logger = logging.getLogger(__name__)

_ADDIN_MARKERS = ("think-cell", "thinkcell")

_REMOVABLE_REL_TYPES = frozenset({
    "oleObject", "package", "tags", "tag", "vmlDrawing",
})

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _strip_addin_shapes(prs: Presentation) -> int:
    """Remove add-in data shapes from all slides.  Returns count removed."""
    removed = 0
    for slide in prs.slides:
        for shape in list(slide.shapes):
            name = (shape.name or "").lower()
            if any(m in name for m in _ADDIN_MARKERS):
                shape.element.getparent().remove(shape.element)
                removed += 1
                continue
            xml_str = etree.tostring(shape.element, encoding="unicode").lower()
            if any(m in xml_str for m in _ADDIN_MARKERS):
                shape.element.getparent().remove(shape.element)
                removed += 1
    return removed


def _prune_orphaned_rels(raw: bytes) -> bytes:
    """Rewrite the PPTX ZIP, removing slide rels whose rIds are unused."""
    src_zip = zipfile.ZipFile(BytesIO(raw), "r")
    files_to_drop: Set[str] = set()
    rels_rewrites: dict[str, bytes] = {}

    for entry in src_zip.namelist():
        if not (entry.startswith("ppt/slides/slide")
                and entry.endswith(".xml")
                and "/_rels/" not in entry):
            continue

        sroot = ET.fromstring(src_zip.read(entry))
        used_rids: Set[str] = set()
        for el in sroot.iter():
            for val in el.attrib.values():
                if val.startswith("rId"):
                    used_rids.add(val)

        rels_path = entry.replace(
            "ppt/slides/", "ppt/slides/_rels/") + ".rels"
        if rels_path not in src_zip.namelist():
            continue

        rroot = ET.fromstring(src_zip.read(rels_path))
        changed = False
        for rel in list(rroot):
            rid = rel.get("Id", "")
            rtype = (rel.get("Type") or "").split("/")[-1]
            target = rel.get("Target", "")
            if rid in used_rids:
                continue
            if rtype not in _REMOVABLE_REL_TYPES:
                continue
            rroot.remove(rel)
            changed = True
            resolved = posixpath.normpath(
                posixpath.join(posixpath.dirname(entry), target))
            files_to_drop.add(resolved)

        if changed:
            ET.register_namespace("", _R_NS)
            rels_rewrites[rels_path] = ET.tostring(
                rroot, xml_declaration=True, encoding="UTF-8")

    if not rels_rewrites and not files_to_drop:
        src_zip.close()
        return raw

    ct_path = "[Content_Types].xml"
    ct_root = ET.fromstring(src_zip.read(ct_path))
    ct_ns = ct_root.tag.split("}")[0] + "}" if "}" in ct_root.tag else ""
    for override in list(ct_root):
        pn = (override.get("PartName") or "").lstrip("/")
        if pn in files_to_drop:
            ct_root.remove(override)

    if ct_ns:
        ET.register_namespace("", ct_ns.strip("{}"))
    ct_bytes = ET.tostring(ct_root, xml_declaration=True, encoding="UTF-8")

    out = BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in src_zip.infolist():
            if item.filename in files_to_drop:
                continue
            if item.filename in rels_rewrites:
                zout.writestr(item, rels_rewrites[item.filename])
            elif item.filename == ct_path:
                zout.writestr(item, ct_bytes)
            else:
                zout.writestr(item, src_zip.read(item.filename))

    src_zip.close()
    out.seek(0)
    return out.read()


def sanitize_template(raw: bytes) -> bytes:
    """Clean a PPTX template: strip add-in shapes, prune orphaned parts.

    Returns the sanitized PPTX bytes ready for storage and later filling.
    """
    prs = Presentation(BytesIO(raw))
    removed = _strip_addin_shapes(prs)

    buf = BytesIO()
    prs.save(buf)
    sanitized = _prune_orphaned_rels(buf.getvalue())

    if removed:
        logger.info("Template sanitized: removed %d add-in shapes", removed)

    return sanitized
