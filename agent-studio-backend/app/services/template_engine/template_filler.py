"""Generic PPTX template filler.

Given a PPTX template with ``{{ }}``, ``{{* }}``, ``{{+ }}``, and
``{{# }}``/``{{/ }}`` placeholders, fill it from a flat or nested JSON/dict
data source and return the populated presentation as bytes.

Slide assembly uses a "build from scratch" pattern: the template is parsed
to determine which slides are needed, then a fresh presentation is assembled
by copying only the required slides.  This avoids python-pptx's broken
slide-cloning/removal internals (orphaned OPC parts, partname collisions).
"""
from __future__ import annotations

import logging
import re
from copy import deepcopy
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx_slide_copier import SlideCopier

from .placeholder_parser import (
    PlaceholderKind,
    detect_loops,
    detect_variants,
    extract_placeholders,
)
from .run_merger import (
    merge_multiline_placeholders,
    paragraph_text,
    replace_all_in_paragraph,
)

logger = logging.getLogger(__name__)

_RE_PLACEHOLDER = re.compile(
    r"\{\{\s*(?P<prefix>[#*/+@]?)\s*"
    r"(?P<path>[a-zA-Z_][a-zA-Z0-9_.]*)"
    r"\s*(?:\|[^}]*)?"
    r"\s*\}\}"
)

_RE_REPEAT = re.compile(
    r"\{\{\+\s*(?P<path>[a-zA-Z_][a-zA-Z0-9_.]*)"
    r"\s*(?:\|[^}]*)?\s*\}\}"
)

_RE_BULLET = re.compile(
    r"\{\{\*\s*(?P<path>[a-zA-Z_][a-zA-Z0-9_.]*)"
    r"\s*(?:\|[^}]*)?\s*\}\}"
)

_RE_VARIANT_MARKER = re.compile(
    r"\{\{\s*@\s*[a-zA-Z_][a-zA-Z0-9_.]*"
    r"\s*(?:\|[^}]*)?\s*\}\}"
)


# =========================================================================
# Data resolution
# =========================================================================

def _resolve(data: Dict[str, Any], path: str) -> Any:
    """Walk a dotted *path* into *data*, returning ``None`` on miss.

    Supports numeric segments as list indices (e.g. ``people.0.name``
    resolves ``data["people"][0]["name"]``).
    """
    current: Any = data
    for part in path.split("."):
        if isinstance(current, dict):
            current = current.get(part)
        elif isinstance(current, list) and part.isdigit():
            idx = int(part)
            current = current[idx] if idx < len(current) else None
        else:
            return None
        if current is None:
            return None
    return current


# =========================================================================
# Text-frame helpers
# =========================================================================

def _iter_text_frames(slide):
    def _walk(shapes):
        for shape in shapes:
            if shape.has_text_frame:
                yield shape.text_frame
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
            if hasattr(shape, "shapes"):
                yield from _walk(shape.shapes)

    yield from _walk(slide.shapes)


def _fill_bullets(paragraph, array_data: List[str]) -> None:
    """Replace a ``{{* }}`` paragraph with one paragraph per list item.

    Clones the formatting of the original paragraph for each new bullet.
    """
    if not array_data:
        paragraph.text = ""
        return

    parent_element = paragraph._p.getparent()
    ref_p_xml = deepcopy(paragraph._p)

    for run_el in ref_p_xml.findall(qn("a:r")):
        for t_el in run_el.findall(qn("a:t")):
            t_el.text = ""

    insert_after = paragraph._p

    for i, bullet_text in enumerate(array_data):
        if i == 0:
            for run in paragraph.runs:
                run.text = ""
            if paragraph.runs:
                paragraph.runs[0].text = bullet_text
            else:
                paragraph.text = bullet_text
            continue

        new_p = deepcopy(ref_p_xml)
        runs = new_p.findall(qn("a:r"))
        if runs:
            t_elements = runs[0].findall(qn("a:t"))
            if t_elements:
                t_elements[0].text = bullet_text
            else:
                t_el = etree.SubElement(runs[0], qn("a:t"))
                t_el.text = bullet_text
        else:
            r_el = etree.SubElement(new_p, qn("a:r"))
            rPr = ref_p_xml.find(qn("a:r"))
            if rPr is not None:
                rPr_copy = rPr.find(qn("a:rPr"))
                if rPr_copy is not None:
                    r_el.insert(0, deepcopy(rPr_copy))
            t_el = etree.SubElement(r_el, qn("a:t"))
            t_el.text = bullet_text

        insert_after.addnext(new_p)
        insert_after = new_p


def _detect_repeat_spans(
    tf,
) -> List[Tuple[str, int, int]]:
    """Find contiguous paragraph spans belonging to the same repeat group.

    Returns ``(array_root, start_idx, end_idx)`` tuples (inclusive indices
    into *tf.paragraphs*).  A span starts at a ``{{+ }}`` paragraph and
    extends through any immediately following ``{{* }}`` paragraph whose
    path shares the same array root.
    """
    paras = list(tf.paragraphs)
    spans: List[Tuple[str, int, int]] = []
    i = 0
    while i < len(paras):
        text = paragraph_text(paras[i])
        m = _RE_REPEAT.search(text)
        if not m:
            i += 1
            continue

        repeat_path = m.group("path")
        parts = repeat_path.rsplit(".", 1)
        if len(parts) != 2:
            i += 1
            continue
        array_root = parts[0]

        span_end = i
        j = i + 1
        while j < len(paras):
            jtext = paragraph_text(paras[j])
            bm = _RE_BULLET.search(jtext)
            if bm:
                bpath = bm.group("path")
                bparts = bpath.rsplit(".", 1)
                if len(bparts) == 2 and bparts[0] == array_root:
                    span_end = j
                    j += 1
                    continue
            rm = _RE_REPEAT.search(jtext)
            if rm:
                rpath = rm.group("path")
                rparts = rpath.rsplit(".", 1)
                if len(rparts) == 2 and rparts[0] == array_root:
                    span_end = j
                    j += 1
                    continue
            break

        spans.append((array_root, i, span_end))
        i = span_end + 1

    return spans


def _set_paragraph_text(p_xml, text: str) -> None:
    """Set the text of the first run in a cloned ``<a:p>`` element."""
    runs = p_xml.findall(qn("a:r"))
    if runs:
        for r in runs:
            for t_el in r.findall(qn("a:t")):
                t_el.text = ""
        t_elements = runs[0].findall(qn("a:t"))
        if t_elements:
            t_elements[0].text = text
        else:
            t_el = etree.SubElement(runs[0], qn("a:t"))
            t_el.text = text
    else:
        r_el = etree.SubElement(p_xml, qn("a:r"))
        t_el = etree.SubElement(r_el, qn("a:t"))
        t_el.text = text


def _expand_bullet_xml(
    ref_p_xml, bullets: List[str], parent, insert_after,
):
    """Clone *ref_p_xml* once per bullet string; return last inserted element."""
    if not bullets:
        new_p = deepcopy(ref_p_xml)
        _set_paragraph_text(new_p, "")
        insert_after.addnext(new_p)
        return new_p

    for bullet_text in bullets:
        new_p = deepcopy(ref_p_xml)
        _set_paragraph_text(new_p, bullet_text)
        insert_after.addnext(new_p)
        insert_after = new_p
    return insert_after


def _fill_repeat_group(
    tf,
    array_root: str,
    span_start: int,
    span_end: int,
    data: Dict[str, Any],
) -> None:
    """Expand a repeat group for each item in the resolved array.

    Clones the template paragraphs (from *span_start* to *span_end*) once
    per array item, filling ``{{+ }}`` fields as text and ``{{* }}`` fields
    as expanded bullet lists.  The original template paragraphs are removed
    after expansion.
    """
    paras = list(tf.paragraphs)
    template_p_xmls = [deepcopy(paras[k]._p) for k in range(span_start, span_end + 1)]

    array_data = _resolve(data, array_root)
    if isinstance(array_data, dict):
        array_data = [array_data]
    elif not isinstance(array_data, list):
        array_data = []

    txBody = paras[0]._p.getparent()

    if span_end + 1 < len(paras):
        insert_before = paras[span_end + 1]._p
    else:
        insert_before = None

    for item in array_data:
        if not isinstance(item, dict):
            item = {}

        for tmpl_p in template_p_xmls:
            raw_text = "".join(
                t.text or ""
                for t in tmpl_p.iter(qn("a:t"))
            )

            bm = _RE_BULLET.search(raw_text)
            if bm:
                bpath = bm.group("path")
                leaf = bpath.rsplit(".", 1)[-1]
                bullets = item.get(leaf, [])
                if not isinstance(bullets, list):
                    bullets = [str(bullets)] if bullets else []
                ref_p = deepcopy(tmpl_p)
                for r in ref_p.findall(qn("a:r")):
                    for t_el in r.findall(qn("a:t")):
                        t_el.text = ""

                for bi, bt in enumerate(bullets):
                    new_p = deepcopy(ref_p)
                    _set_paragraph_text(new_p, bt)
                    if insert_before is not None:
                        insert_before.addprevious(new_p)
                    else:
                        txBody.append(new_p)
                if not bullets:
                    empty_p = deepcopy(ref_p)
                    _set_paragraph_text(empty_p, "")
                    if insert_before is not None:
                        insert_before.addprevious(empty_p)
                    else:
                        txBody.append(empty_p)
                continue

            rm = _RE_REPEAT.search(raw_text)
            if rm:
                rpath = rm.group("path")
                leaf = rpath.rsplit(".", 1)[-1]
                value = item.get(leaf, "")
                if value is None:
                    value = ""
                new_p = deepcopy(tmpl_p)
                full_tag = rm.group(0)
                combined = "".join(
                    t.text or ""
                    for t in new_p.iter(qn("a:t"))
                )
                if full_tag in combined:
                    _replace_in_p_xml(new_p, full_tag, str(value))
                else:
                    _set_paragraph_text(new_p, str(value))
                if insert_before is not None:
                    insert_before.addprevious(new_p)
                else:
                    txBody.append(new_p)
                continue

            new_p = deepcopy(tmpl_p)
            if insert_before is not None:
                insert_before.addprevious(new_p)
            else:
                txBody.append(new_p)

    if not array_data:
        for tmpl_p in template_p_xmls:
            empty_p = deepcopy(tmpl_p)
            for r in empty_p.findall(qn("a:r")):
                for t_el in r.findall(qn("a:t")):
                    t_el.text = ""
            if insert_before is not None:
                insert_before.addprevious(empty_p)
            else:
                txBody.append(empty_p)

    for k in range(span_start, span_end + 1):
        txBody.remove(paras[k]._p)

    logger.debug(
        "Expanded repeat group %s: %d items, %d template paragraphs",
        array_root, len(array_data), len(template_p_xmls),
    )


def _replace_in_p_xml(p_xml, old: str, new: str) -> None:
    """Replace *old* with *new* across runs in a raw ``<a:p>`` element."""
    runs = p_xml.findall(qn("a:r"))
    if not runs:
        return

    texts = []
    for r in runs:
        for t in r.findall(qn("a:t")):
            texts.append((r, t))

    full = "".join(t.text or "" for _, t in texts)
    idx = full.find(old)
    if idx == -1:
        return

    boundaries = []
    pos = 0
    for r, t in texts:
        end = pos + len(t.text or "")
        boundaries.append((pos, end, r, t))
        pos = end

    match_start = idx
    match_end = idx + len(old)

    first_i = None
    for i, (s, e, _, _) in enumerate(boundaries):
        if s <= match_start < e:
            first_i = i
            break

    last_i = None
    for i, (s, e, _, _) in enumerate(boundaries):
        if s < match_end <= e:
            last_i = i
            break

    if first_i is None or last_i is None:
        return

    if first_i == last_i:
        _, _, _, t_el = boundaries[first_i]
        local_s = match_start - boundaries[first_i][0]
        local_e = match_end - boundaries[first_i][0]
        t_el.text = (t_el.text or "")[:local_s] + new + (t_el.text or "")[local_e:]
    else:
        _, _, _, ft = boundaries[first_i]
        local_s = match_start - boundaries[first_i][0]
        ft.text = (ft.text or "")[:local_s] + new

        for mid in range(first_i + 1, last_i):
            boundaries[mid][3].text = ""

        _, _, _, lt = boundaries[last_i]
        local_e = match_end - boundaries[last_i][0]
        lt.text = (lt.text or "")[local_e:]


# =========================================================================
# Per-slide filling
# =========================================================================

def _fill_slide(slide, data: Dict[str, Any]) -> None:
    """Fill ``{{ }}``, ``{{* }}``, and ``{{+ }}`` placeholders on a slide."""
    for tf in _iter_text_frames(slide):
        merge_multiline_placeholders(tf)

        spans = _detect_repeat_spans(tf)
        if spans:
            for array_root, s_start, s_end in reversed(spans):
                _fill_repeat_group(tf, array_root, s_start, s_end, data)

    for tf in _iter_text_frames(slide):
        paragraphs = list(tf.paragraphs)
        for para in paragraphs:
            full_text = paragraph_text(para)

            if _RE_REPEAT.search(full_text):
                continue

            bullet_match = _RE_BULLET.search(full_text)
            if bullet_match:
                path = bullet_match.group("path")
                parts = path.rsplit(".", 1)
                if len(parts) == 2:
                    maybe_array = _resolve(data, parts[0])
                    if isinstance(maybe_array, list):
                        continue
                arr = _resolve(data, path)
                if not isinstance(arr, list):
                    arr = [str(arr)] if arr else []
                _fill_bullets(para, arr)
                continue

            for m in _RE_PLACEHOLDER.finditer(full_text):
                prefix = m.group("prefix")
                if prefix in ("#", "/", "*", "+", "@"):
                    continue
                path = m.group("path")
                value = _resolve(data, path)
                if value is None:
                    value = ""
                replace_all_in_paragraph(para, m.group(0), str(value))


def _remove_loop_markers(slide) -> None:
    """Strip ``{{# }}`` and ``{{/ }}`` marker text from a slide."""
    for tf in _iter_text_frames(slide):
        merge_multiline_placeholders(tf)
        for para in tf.paragraphs:
            full_text = paragraph_text(para)
            for m in re.finditer(
                r"\{\{\s*[#/]\s*[a-zA-Z_][a-zA-Z0-9_.]*"
                r"\s*(?:\|[^}]*)?\s*\}\}",
                full_text,
            ):
                replace_all_in_paragraph(para, m.group(0), "")


def _strip_variant_markers(slide) -> None:
    """Remove ``{{@ ... }}`` marker text from a slide.

    If a top-level shape contains ONLY variant markers (and whitespace),
    the entire shape is removed so any visual styling (colored boxes,
    borders) the template author used doesn't leak into the output.
    """
    shapes_to_remove = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        merge_multiline_placeholders(tf)

        all_text = "".join(paragraph_text(p) for p in tf.paragraphs)
        if _RE_VARIANT_MARKER.search(all_text):
            text_without_markers = _RE_VARIANT_MARKER.sub("", all_text).strip()
            if not text_without_markers:
                shapes_to_remove.append(shape)
                continue

        for para in tf.paragraphs:
            full_text = paragraph_text(para)
            for m in _RE_VARIANT_MARKER.finditer(full_text):
                replace_all_in_paragraph(para, m.group(0), "")

    for shape in shapes_to_remove:
        shape._element.getparent().remove(shape._element)
        logger.debug("Removed variant-marker-only shape from slide")


# =========================================================================
# Variant selection (pure logic)
# =========================================================================

def _select_variant(available: Dict[int, List[int]], count: int) -> int:
    """Pick the best variant count for the given *count*.

    Exact match first, then nearest lower-or-equal, then nearest higher,
    and finally the maximum available as a last resort.
    """
    counts = sorted(available.keys())
    if count in available:
        return count
    if count > counts[-1]:
        return counts[-1]
    if count < counts[0]:
        return counts[0]
    return max(c for c in counts if c <= count)


# =========================================================================
# Slide plan builder (pure logic -- no I/O, no mutations)
# =========================================================================

def _build_slide_plan(
    prs: Presentation,
    placeholders: List,
    variant_groups: List[Dict],
    loops: List[Dict],
    data: Dict[str, Any],
) -> List[Tuple[int, Dict[str, Any]]]:
    """Compute an ordered list of ``(source_slide_index, data_dict)`` tuples.

    Each tuple describes one slide that should appear in the final output.
    """
    total = len(prs.slides)

    # Identify slides consumed by loop markers (these are never output)
    loop_marker_indices: set = set()
    loop_body_ranges: List[Tuple[int, int, str]] = []
    for loop in loops:
        loop_marker_indices.add(loop["start_slide"])
        loop_marker_indices.add(loop["end_slide"])
        body_start = loop["start_slide"] + 1
        body_end = loop["end_slide"] - 1
        if body_end >= body_start:
            loop_body_ranges.append((body_start, body_end, loop["name"]))

    loop_body_indices: set = set()
    for bs, be, _ in loop_body_ranges:
        loop_body_indices.update(range(bs, be + 1))

    # Identify variant slides: determine which to keep, which to discard
    variant_keep: set = set()
    variant_discard: set = set()
    for group in variant_groups:
        array_data = _resolve(data, group["name"])
        count = len(array_data) if isinstance(array_data, list) else 0
        selected = _select_variant(group["variants"], count)
        for variant_count, slide_indices in group["variants"].items():
            if variant_count == selected:
                variant_keep.update(slide_indices)
            else:
                variant_discard.update(slide_indices)

    plan: List[Tuple[int, Dict[str, Any]]] = []

    for idx in range(total):
        if idx in variant_discard:
            continue
        if idx in loop_marker_indices:
            continue

        # Loop body slides get repeated per array item
        in_loop = False
        for bs, be, loop_name in loop_body_ranges:
            if bs <= idx <= be:
                in_loop = True
                if idx == bs:
                    array_data = data.get(loop_name, [])
                    if not isinstance(array_data, list):
                        array_data = []
                    for item in array_data:
                        for body_idx in range(bs, be + 1):
                            scoped = {"item": item}
                            scoped.update(data)
                            plan.append((body_idx, scoped))
                break

        if in_loop:
            continue

        plan.append((idx, data))

    return plan


# =========================================================================
# Presentation assembly
# =========================================================================

def _remove_all_slides(prs: Presentation) -> None:
    """Remove every slide from *prs*, leaving themes and layouts intact."""
    while len(prs.slides._sldIdLst):
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


# =========================================================================
# Public API
# =========================================================================

def fill_template(template_path: str, data: Dict[str, Any]) -> bytes:
    """Fill *template_path* with *data* and return the PPTX as bytes.

    Uses a "build from scratch" pattern: the template is opened read-only
    as a source, a second copy is opened as the destination (preserving
    themes and layouts), all slides are removed from the destination, and
    only the required slides are copied over and filled.
    """
    src = Presentation(template_path)
    placeholders = extract_placeholders(prs=src)
    variant_groups = detect_variants(placeholders)
    loops = detect_loops(placeholders)

    slide_plan = _build_slide_plan(src, placeholders, variant_groups, loops, data)

    logger.info(
        "fill_template: %d source slides, plan has %d output slides",
        len(src.slides), len(slide_plan),
    )

    dst = Presentation(template_path)
    _remove_all_slides(dst)

    for src_idx, _ in slide_plan:
        SlideCopier.copy_slide(src, src_idx, dst)

    for slide, (_, slide_data) in zip(dst.slides, slide_plan):
        _strip_variant_markers(slide)
        _remove_loop_markers(slide)
        _fill_slide(slide, slide_data)

    buf = BytesIO()
    dst.save(buf)
    return buf.getvalue()


def fill_template_to_file(template_path: str, data: Dict[str, Any],
                          output_path: str) -> str:
    """Fill and write directly to *output_path*.  Returns the path."""
    content = fill_template(template_path, data)
    with open(output_path, "wb") as f:
        f.write(content)
    return output_path
